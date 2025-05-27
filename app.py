import os
import streamlit as st
import tempfile
import logging
import requests
import base64
import time
import pandas as pd
import numpy as np
from pptx import Presentation
from docx import Document as DocxDocument
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.callbacks import StdOutCallbackHandler

# Fix for LangSmithTracer import error
try:
    from langchain.callbacks.tracers.langchain import LangChainTracer
    has_langsmith = True
except ImportError:
    has_langsmith = False
    class LangChainTracer:
        def __init__(self, *args, **kwargs):
            pass

# Updated Pinecone import for version compatibility
import pinecone
from langchain_pinecone import PineconeVectorStore
import google.generativeai as genai

# Import speech recognition for voice input
try:
    import speech_recognition as sr
    has_speech_recognition = True
except ImportError:
    has_speech_recognition = False

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("langchain_processes.log"),
        logging.StreamHandler()
    ]
)

logger = logging.getLogger("pdf_chatbot")

# Voice functionality using ElevenLabs
class VoiceHandler:
    def __init__(self):
        self.api_key = st.secrets.get("ELEVENLABS_API_KEY", "")
        self.voice_id = st.secrets.get("ELEVENLABS_VOICE_ID", "21m00Tcm4TlvDq8ikWAM")  # Default voice
        self.base_url = "https://api.elevenlabs.io/v1"
        
    def text_to_speech(self, text):
        """Convert text to speech using ElevenLabs API"""
        if not self.api_key:
            logger.warning("ElevenLabs API key not found")
            return None
            
        url = f"{self.base_url}/text-to-speech/{self.voice_id}"
        
        headers = {
            "Accept": "audio/mpeg",
            "Content-Type": "application/json",
            "xi-api-key": self.api_key
        }
        
        data = {
            "text": text,
            "model_id": "eleven_monolingual_v1",
            "voice_settings": {
                "stability": 0.5,
                "similarity_boost": 0.5
            }
        }
        
        try:
            response = requests.post(url, json=data, headers=headers)
            if response.status_code == 200:
                return response.content
            else:
                logger.error(f"ElevenLabs API error: {response.status_code} - {response.text}")
                return None
        except Exception as e:
            logger.error(f"Error calling ElevenLabs API: {str(e)}")
            return None
    
    def get_available_voices(self):
        """Get list of available voices from ElevenLabs"""
        if not self.api_key:
            return []
            
        url = f"{self.base_url}/voices"
        headers = {"xi-api-key": self.api_key}
        
        try:
            response = requests.get(url, headers=headers)
            if response.status_code == 200:
                voices = response.json().get("voices", [])
                return [(voice["voice_id"], voice["name"]) for voice in voices]
            return []
        except Exception as e:
            logger.error(f"Error fetching voices: {str(e)}")
            return []

class DIDHandler:
    def __init__(self):
        self.api_key = st.secrets.get("DID_API_KEY", "")
        self.base_url = "https://api.d-id.com"
        # Use the provided presenter ID from secrets or fall back to default
        self.presenter_id = st.secrets.get("DID_PRESENTER_ID", "amy-jcwCkr1grs")  # Amy is the default presenter
        
    def get_available_presenters(self):
        """Get list of available presenters from D-ID"""
        if not self.api_key:
            return []
            
        url = f"{self.base_url}/presenters"
        headers = {
            "Authorization": f"Basic {self.api_key}",
            "Content-Type": "application/json"
        }
        
        try:
            response = requests.get(url, headers=headers)
            if response.status_code == 200:
                presenters = response.json().get("presenters", [])
                return [(presenter["id"], presenter["name"]) for presenter in presenters]
            return []
        except Exception as e:
            logger.error(f"Error fetching presenters: {str(e)}")
            return []
    
    def upload_custom_image(self, image_file):
        """Upload a custom image to D-ID and return the image URL"""
        if not self.api_key:
            logger.warning("D-ID API key not found")
            return None
            
        url = f"{self.base_url}/images"
        
        headers = {
            "Authorization": f"Basic {self.api_key}"
        }
        
        try:
            # Prepare the file for upload
            files = {
                'image': (image_file.name, image_file.getvalue(), image_file.type)
            }
            
            response = requests.post(url, headers=headers, files=files)
            
            if response.status_code == 201:
                result = response.json()
                image_url = result.get("url")
                logger.info(f"Successfully uploaded custom image: {image_url}")
                return image_url
            else:
                logger.error(f"D-ID image upload error: {response.status_code} - {response.text}")
                return None
                
        except Exception as e:
            logger.error(f"Error uploading image to D-ID: {str(e)}")
            return None
    
    def create_talk_with_custom_image(self, text, image_url):
        """Create a talking video using D-ID API with custom image"""
        if not self.api_key:
            logger.warning("D-ID API key not found")
            return None
            
        url = f"{self.base_url}/talks"
        
        headers = {
            "Authorization": f"Basic {self.api_key}",
            "Content-Type": "application/json"
        }
        
        data = {
            "script": {
                "type": "text",
                "input": text
            },
            "source_url": image_url,  # Use custom image instead of presenter_id
            "driver_id": "wav2lip"
        }
        
        try:
            response = requests.post(url, json=data, headers=headers)
            if response.status_code == 201:
                return response.json().get("id")
            else:
                logger.error(f"D-ID API error: {response.status_code} - {response.text}")
                return None
        except Exception as e:
            logger.error(f"Error calling D-ID API: {str(e)}")
            return None
            
    def create_talk(self, text, custom_image_url=None):
        """Create a talking video using D-ID API with either custom image or presenter"""
        if custom_image_url:
            return self.create_talk_with_custom_image(text, custom_image_url)
        
        # Original method using presenter_id
        if not self.api_key or not self.presenter_id:
            logger.warning("D-ID API key or presenter ID not found")
            return None
            
        url = f"{self.base_url}/talks"
        
        headers = {
            "Authorization": f"Basic {self.api_key}",
            "Content-Type": "application/json"
        }
        
        data = {
            "script": {
                "type": "text",
                "input": text
            },
            "presenter_id": self.presenter_id,
            "driver_id": "wav2lip"
        }
        
        try:
            response = requests.post(url, json=data, headers=headers)
            if response.status_code == 201:
                return response.json().get("id")
            else:
                logger.error(f"D-ID API error: {response.status_code} - {response.text}")
                return None
        except Exception as e:
            logger.error(f"Error calling D-ID API: {str(e)}")
            return None
            
    def get_talk_url(self, talk_id):
        """Get the video URL for a completed talk"""
        if not self.api_key:
            return None
            
        url = f"{self.base_url}/talks/{talk_id}"
        headers = {"Authorization": f"Basic {self.api_key}"}
        
        max_attempts = 10
        attempt = 0
        
        while attempt < max_attempts:
            try:
                response = requests.get(url, headers=headers)
                if response.status_code == 200:
                    result = response.json()
                    status = result.get("status")
                    
                    if status == "done":
                        return result.get("result_url")
                    elif status in ["error", "failed"]:
                        logger.error(f"Talk generation failed: {result.get('error', 'Unknown error')}")
                        return None
                        
                time.sleep(2)  # Wait before checking again
                attempt += 1
            except Exception as e:
                logger.error(f"Error checking talk status: {str(e)}")
                return None
                
        logger.error("Timeout waiting for talk to complete")
        return None

# Set up environment variables from Streamlit secrets
def setup_environment():
    # Set required API keys
    os.environ["PINECONE_API_KEY"] = st.secrets["PINECONE_API_KEY"]
    os.environ["GOOGLE_API_KEY"] = st.secrets["GOOGLE_API_KEY"]
    
    # Set Pinecone environment for v2 compatibility
    if "PINECONE_ENVIRONMENT" in st.secrets:
        os.environ["PINECONE_ENVIRONMENT"] = st.secrets["PINECONE_ENVIRONMENT"]

    # Set LangSmith variables if they exist in secrets
    if "LANGCHAIN_API_KEY" in st.secrets:
        os.environ["LANGCHAIN_API_KEY"] = st.secrets["LANGCHAIN_API_KEY"]
        os.environ["LANGCHAIN_PROJECT"] = st.secrets.get("LANGCHAIN_PROJECT", "pdf-chatbot")
        os.environ["LANGCHAIN_TRACING_V2"] = st.secrets.get("LANGCHAIN_TRACING_V2", "true")
        logger.info("LangSmith tracing enabled")
    else:
        logger.info("LangSmith tracing not configured")

# Initialize environment
setup_environment()

# Initialize Pinecone with version detection
try:
    pinecone_version = pinecone.__version__.split('.')[0]
    logger.info(f"Detected Pinecone SDK version: {pinecone.__version__}")
    
    if int(pinecone_version) >= 4:
        pc = pinecone.Pinecone(api_key=os.environ["PINECONE_API_KEY"])
        logger.info("Initialized Pinecone v4+ client")
    else:
        pinecone.init(
            api_key=os.environ["PINECONE_API_KEY"],
            environment=os.environ.get("PINECONE_ENVIRONMENT", "gcp-starter")
        )
        pc = pinecone
        logger.info("Initialized Pinecone v2-v3 client")
except Exception as e:
    logger.error(f"Error initializing Pinecone: {str(e)}")
    st.error(f"Error initializing Pinecone: {str(e)}")
    pc = None

# Initialize Gemini
genai.configure(api_key=os.environ["GOOGLE_API_KEY"])

class PDFChatbot:
    def __init__(self):
        logger.info("Initializing PDF Chatbot")

        self.embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        self.index_name = "pdf-ai"
        self.pc = pc

        callbacks = None
        if "LANGCHAIN_API_KEY" in os.environ and has_langsmith:
            try:
                tracer = LangChainTracer(
                    project_name=os.environ.get("LANGCHAIN_PROJECT", "pdf-chatbot")
                )
                callbacks = [StdOutCallbackHandler(), tracer]
                logger.info("LangSmith tracing enabled with callback handlers")
            except Exception as e:
                logger.warning(f"Failed to initialize LangSmith tracing: {e}")
                callbacks = None

        try:
            self.llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")
            logger.info("Successfully initialized LLM with gemini-1.5-flash model")
        except Exception as e:
            logger.warning(f"Error initializing gemini-1.5-flash: {e}")
            try:
                self.llm = ChatGoogleGenerativeAI(model="gemini-1.0-pro")
                logger.info("Successfully initialized LLM with gemini-1.0-pro model")
            except Exception as e2:
                logger.warning(f"Error initializing gemini-1.0-pro: {e2}")
                self.llm = ChatGoogleGenerativeAI(model="models/gemini-pro")
                logger.info("Trying legacy model name format: models/gemini-pro")

        self._setup_pinecone_index()

        self.memory = ConversationBufferMemory(
            memory_key="chat_history",
            return_messages=True,
            output_key="answer"
        )

        self.callbacks = callbacks
        logger.info("PDF Chatbot initialization complete")

    def _setup_pinecone_index(self):
        if self.pc is None:
            logger.error("Pinecone client not initialized")
            st.error("Pinecone client not initialized")
            return
            
        try:
            pinecone_version = pinecone.__version__.split('.')[0]
            
            if int(pinecone_version) >= 4:
                index_names = [idx["name"] for idx in self.pc.list_indexes()]
                if self.index_name not in index_names:
                    logger.info(f"Creating new Pinecone index: {self.index_name}")
                    try:
                        self.pc.create_index(
                            name=self.index_name,
                            dimension=384,
                            metric="cosine"
                        )
                        logger.info(f"Created new Pinecone index: {self.index_name}")
                    except Exception as e:
                        logger.error(f"Failed to create Pinecone index: {str(e)}")
                        st.error(f"Failed to create Pinecone index: {str(e)}")
                else:
                    logger.info(f"Using existing Pinecone index: {self.index_name}")
            else:
                index_list = self.pc.list_indexes()
                if self.index_name not in index_list:
                    logger.info(f"Creating new Pinecone index: {self.index_name}")
                    self.pc.create_index(
                        name=self.index_name,
                        dimension=384,
                        metric="cosine"
                    )
                    logger.info(f"Created new Pinecone index: {self.index_name}")
                else:
                    logger.info(f"Using existing Pinecone index: {self.index_name}")
        except Exception as e:
            logger.error(f"Error with Pinecone index setup: {e}")
            st.error(f"Error with Pinecone index setup: {e}")

    def extract_text_from_pdf(self, pdf_file):
        logger.info(f"Starting text extraction from PDF: {pdf_file.name}")

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            tmp_file.write(pdf_file.getvalue())
            pdf_path = tmp_file.name

        logger.info(f"Temporary PDF file created: {pdf_path}")

        try:
            loader = PyMuPDFLoader(pdf_path)
            documents = loader.load()
            logger.info(f"Extracted {len(documents)} pages from PDF")
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            st.error(f"Error extracting text: {str(e)}")
            documents = []
        finally:
            os.unlink(pdf_path)
            logger.info(f"Temporary PDF file deleted")

        return documents

    def extract_text_from_txt(self, txt_file):
        """Extract text from TXT file and return as Document objects"""
        logger.info(f"Starting text extraction from TXT: {txt_file.name}")
        
        try:
            # Read the text file content
            text_content = txt_file.getvalue().decode('utf-8')
            logger.info(f"Successfully read TXT file with {len(text_content)} characters")
            
            # Create a Document object using LangChain's Document class
            from langchain_core.documents import Document
            document = Document(
                page_content=text_content,
                metadata={
                    "source": txt_file.name,
                    "file_type": "txt",
                    "page": 1  # TXT files are treated as single page
                }
            )
            
            logger.info("Created Document object from TXT file")
            return [document]  # Return as list to match PDF structure
            
        except UnicodeDecodeError:
            logger.error("Error: Unable to decode TXT file. Please ensure it's UTF-8 encoded.")
            st.error("Error: Unable to decode TXT file. Please ensure it's UTF-8 encoded.")
            return []
        except Exception as e:
            logger.error(f"Error extracting text from TXT: {str(e)}")
            st.error(f"Error extracting text: {str(e)}")
            return []

    def extract_text_from_docx(self, docx_file):
        """Extract text from DOCX file and return as Document objects"""
        logger.info(f"Starting text extraction from DOCX: {docx_file.name}")
        
        try:
            # Create a temporary file for the DOCX
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                tmp_file.write(docx_file.getvalue())
                docx_path = tmp_file.name
            
            logger.info(f"Temporary DOCX file created: {docx_path}")
            
            try:
                # Load the DOCX document
                doc = DocxDocument(docx_path)
                
                # Extract text from all paragraphs
                text_content = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():  # Only add non-empty paragraphs
                        text_content.append(paragraph.text.strip())
                
                # Extract text from tables if present
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_content.append(cell.text.strip())
                
                # Join all text with newlines
                full_text = '\n\n'.join(text_content)
                
                logger.info(f"Successfully extracted DOCX with {len(full_text)} characters from {len(text_content)} text blocks")
                
                # Create a Document object using LangChain's Document class
                from langchain_core.documents import Document
                document = Document(
                    page_content=full_text,
                    metadata={
                        "source": docx_file.name,
                        "file_type": "docx",
                        "page": 1,  # DOCX files are treated as single page
                        "paragraphs": len(doc.paragraphs),
                        "tables": len(doc.tables)
                    }
                )
                
                logger.info("Created Document object from DOCX file")
                return [document]  # Return as list to match PDF structure
                
            except Exception as e:
                logger.error(f"Error processing DOCX content: {str(e)}")
                st.error(f"Error processing DOCX content: {str(e)}")
                return []
            finally:
                # Clean up temporary file
                try:
                    os.unlink(docx_path)
                    logger.info("Temporary DOCX file deleted")
                except Exception as cleanup_error:
                    logger.warning(f"Could not delete temporary DOCX file: {cleanup_error}")
                
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            st.error(f"Error extracting text from DOCX: {str(e)}")
            return []

    def extract_text_from_pptx(self, pptx_file):
        """Extract text from PPTX file and return as Document objects"""
        logger.info(f"Starting text extraction from PPTX: {pptx_file.name}")
        
        try:
            # Create a temporary file for the PPTX
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(pptx_file.getvalue())
                pptx_path = tmp_file.name
            
            logger.info(f"Temporary PPTX file created: {pptx_path}")
            
            try:
                # Load the PPTX presentation
                prs = Presentation(pptx_path)
                
                documents = []
                total_slides = len(prs.slides)
                logger.info(f"Processing {total_slides} slides from PPTX")
                
                # Extract text from each slide
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    
                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                        
                        # Extract text from tables if present
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        slide_text.append(cell.text.strip())
                    
                    # Create document for this slide if it has content
                    if slide_text:
                        slide_content = '\n\n'.join(slide_text)
                        
                        # Create a Document object using LangChain's Document class
                        from langchain_core.documents import Document
                        document = Document(
                            page_content=slide_content,
                            metadata={
                                "source": pptx_file.name,
                                "file_type": "pptx",
                                "slide": slide_num,
                                "total_slides": total_slides,
                                "page": slide_num  # For compatibility with existing code
                            }
                        )
                        documents.append(document)
                        logger.info(f"Processed slide {slide_num} with {len(slide_content)} characters")
                    else:
                        logger.info(f"Slide {slide_num} has no extractable text content")
                
                logger.info(f"Successfully extracted text from {len(documents)} slides out of {total_slides} total slides")
                return documents
                
            except Exception as e:
                logger.error(f"Error processing PPTX content: {str(e)}")
                st.error(f"Error processing PPTX content: {str(e)}")
                return []
            finally:
                # Clean up temporary file
                try:
                    os.unlink(pptx_path)
                    logger.info("Temporary PPTX file deleted")
                except Exception as cleanup_error:
                    logger.warning(f"Could not delete temporary PPTX file: {cleanup_error}")
                    
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            st.error(f"Error extracting text from PPTX: {str(e)}")
            return []

    def extract_text_from_excel(self, excel_file):
        """Extract text from Excel file and return as Document objects"""
        logger.info(f"Starting text extraction from Excel: {excel_file.name}")
        
        try:
            # Create a temporary file for the Excel file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
                tmp_file.write(excel_file.getvalue())
                excel_path = tmp_file.name
            
            logger.info(f"Temporary Excel file created: {excel_path}")
            
            try:
                # Read all sheets from the Excel file
                excel_data = pd.read_excel(excel_path, sheet_name=None, engine='openpyxl')
                
                documents = []
                total_sheets = len(excel_data)
                logger.info(f"Processing {total_sheets} sheets from Excel file")
                
                # Process each sheet
                for sheet_num, (sheet_name, df) in enumerate(excel_data.items(), 1):
                    logger.info(f"Processing sheet: {sheet_name} with {len(df)} rows and {len(df.columns)} columns")
                    
                    # Skip empty sheets
                    if df.empty:
                        logger.info(f"Sheet '{sheet_name}' is empty, skipping")
                        continue
                    
                    # Convert DataFrame to markdown-style table string
                    sheet_content = self._dataframe_to_markdown_table(df, sheet_name)
                    
                    if sheet_content.strip():
                        # Create a Document object using LangChain's Document class
                        from langchain_core.documents import Document
                        document = Document(
                            page_content=sheet_content,
                            metadata={
                                "source": excel_file.name,
                                "file_type": "excel",
                                "sheet_name": sheet_name,
                                "sheet_number": sheet_num,
                                "total_sheets": total_sheets,
                                "rows": len(df),
                                "columns": len(df.columns),
                                "page": sheet_num  # For compatibility with existing code
                            }
                        )
                        documents.append(document)
                        logger.info(f"Processed sheet '{sheet_name}' with {len(sheet_content)} characters")
                    else:
                        logger.info(f"Sheet '{sheet_name}' has no extractable content")
                
                logger.info(f"Successfully extracted text from {len(documents)} sheets out of {total_sheets} total sheets")
                return documents
                
            except Exception as e:
                logger.error(f"Error processing Excel content: {str(e)}")
                st.error(f"Error processing Excel content: {str(e)}")
                return []
            finally:
                # Clean up temporary file
                try:
                    os.unlink(excel_path)
                    logger.info("Temporary Excel file deleted")
                except Exception as cleanup_error:
                    logger.warning(f"Could not delete temporary Excel file: {cleanup_error}")
                    
        except Exception as e:
            logger.error(f"Error extracting text from Excel: {str(e)}")
            st.error(f"Error extracting text from Excel: {str(e)}")
            return []

    def _dataframe_to_markdown_table(self, df, sheet_name):
        """Convert pandas DataFrame to a readable markdown-style table string"""
        try:
            # Start with sheet information
            content_lines = [
                f"## Sheet: {sheet_name}",
                f"**Rows:** {len(df)} | **Columns:** {len(df.columns)}",
                ""
            ]
            
            # Handle completely empty DataFrame
            if df.empty:
                content_lines.append("*This sheet contains no data.*")
                return "\n".join(content_lines)
            
            # Clean column names and handle unnamed columns
            clean_columns = []
            for i, col in enumerate(df.columns):
                if pd.isna(col) or str(col).startswith('Unnamed:'):
                    clean_columns.append(f"Column_{i+1}")
                else:
                    clean_columns.append(str(col).strip())
            
            df.columns = clean_columns
            
            # Limit the number of rows for very large sheets (to prevent memory issues)
            max_rows = 1000
            if len(df) > max_rows:
                df_sample = df.head(max_rows)
                content_lines.append(f"*Note: Showing first {max_rows} rows out of {len(df)} total rows*")
                content_lines.append("")
            else:
                df_sample = df
            
            # Convert DataFrame to string representation
            # Fill NaN values with empty strings for better readability
            df_clean = df_sample.fillna('')
            
            # Convert to string, handling different data types
            df_str = df_clean.astype(str)
            
            # Create table header
            header_row = "| " + " | ".join(df_str.columns) + " |"
            separator_row = "| " + " | ".join(["---"] * len(df_str.columns)) + " |"
            
            content_lines.extend([header_row, separator_row])
            
            # Add data rows
            for _, row in df_str.iterrows():
                # Clean row values and limit length to prevent extremely long cells
                clean_row_values = []
                for val in row:
                    clean_val = str(val).strip()
                    # Limit cell content to 100 characters to prevent excessive length
                    if len(clean_val) > 100:
                        clean_val = clean_val[:97] + "..."
                    # Replace newlines and pipe characters that would break markdown
                    clean_val = clean_val.replace('\n', ' ').replace('\r', ' ').replace('|', '\\|')
                    clean_row_values.append(clean_val)
                
                data_row = "| " + " | ".join(clean_row_values) + " |"
                content_lines.append(data_row)
            
            # Add summary statistics for numeric columns
            numeric_cols = df_sample.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                content_lines.extend([
                    "",
                    "### Summary Statistics for Numeric Columns:",
                    ""
                ])
                
                for col in numeric_cols:
                    if col in df_sample.columns:
                        col_data = pd.to_numeric(df_sample[col], errors='coerce').dropna()
                        if len(col_data) > 0:
                            content_lines.append(f"**{col}:** Count: {len(col_data)}, Mean: {col_data.mean():.2f}, Min: {col_data.min()}, Max: {col_data.max()}")
            
            return "\n".join(content_lines)
            
        except Exception as e:
            logger.error(f"Error converting DataFrame to markdown: {str(e)}")
            return f"## Sheet: {sheet_name}\n\nError processing sheet data: {str(e)}"

    def chunk_text(self, documents):
        logger.info("Starting text chunking process")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )
        chunks = text_splitter.split_documents(documents)

        logger.info(f"Created {len(chunks)} chunks from documents")

        if chunks:
            logger.info(f"Sample chunk content: {chunks[0].page_content[:100]}...")

        return chunks

    def create_embeddings_and_store(self, chunks, namespace):
        logger.info(f"Creating embeddings for {len(chunks)} chunks and storing in Pinecone namespace: {namespace}")

        try:
            pinecone_version = pinecone.__version__.split('.')[0]
            
            if int(pinecone_version) >= 4:
                index = self.pc.Index(self.index_name)
                try:
                    index.delete(namespace=namespace, delete_all=True)
                    logger.info(f"Deleted existing vectors in namespace: {namespace}")
                except Exception as e:
                    logger.warning(f"No existing vectors to delete or error: {str(e)}")
            else:
                index = self.pc.Index(self.index_name)
                try:
                    index.delete(deleteAll=True, namespace=namespace)
                    logger.info(f"Deleted existing vectors in namespace: {namespace}")
                except Exception as e:
                    logger.warning(f"No existing vectors to delete or error: {str(e)}")
                
            vectorstore = PineconeVectorStore.from_documents(
                documents=chunks,
                embedding=self.embeddings,
                index_name=self.index_name,
                namespace=namespace
            )
            logger.info("Embeddings created and stored successfully")
            return vectorstore
        except Exception as e:
            logger.error(f"Error creating embeddings: {str(e)}")
            st.error(f"Error creating embeddings: {str(e)}")
            
            logger.error(f"Error type: {type(e)}")
            logger.error(f"Error details: {str(e)}")
            
            try:
                logger.info("Attempting alternative approach for storing embeddings...")
                vectorstore = PineconeVectorStore.from_documents(
                    documents=chunks,
                    embedding=self.embeddings,
                    index=index,
                    namespace=namespace
                )
                logger.info("Alternative approach successful")
                return vectorstore
            except Exception as e2:
                logger.error(f"Alternative approach also failed: {str(e2)}")
                raise e2

    def get_conversational_chain(self, vectorstore):
        logger.info("Creating conversational retrieval chain")

        retriever = vectorstore.as_retriever(search_kwargs={"k": 5})

        chain = ConversationalRetrievalChain.from_llm(
            llm=self.llm,
            retriever=retriever,
            memory=self.memory,
            return_source_documents=True,
            output_key="answer"
        )

        logger.info("Conversational retrieval chain created")
        return chain

    def process_query(self, chain, query):
        logger.info(f"Processing query: {query}")

        try:
            result = chain.invoke({"question": query})
            logger.info("Used chain.invoke() successfully")
        except AttributeError:
            result = chain({"question": query})
            logger.info("Used chain() call successfully")

        logger.info(f"Retrieved {len(result['source_documents'])} source documents")
        logger.info(f"Generated answer (first 100 chars): {result['answer'][:100]}...")

        return result["answer"], result["source_documents"]


def create_audio_player(audio_content, message_id, should_autoplay=False):
    """Create an HTML audio player for the generated audio with proper playback control"""
    audio_base64 = base64.b64encode(audio_content).decode()
    
    # Only add autoplay attribute, but control it via JavaScript
    audio_html = f"""
    <div id="audio_container_{message_id}">
        <audio id="audio_{message_id}" controls style="width: 100%;">
            <source src="data:audio/mpeg;base64,{audio_base64}" type="audio/mpeg">
            Your browser does not support the audio element.
        </audio>
    </div>
    <script>
        (function() {{
            const audioElement = document.getElementById('audio_{message_id}');
            const containerId = 'audio_container_{message_id}';
            
            // Check if this audio has already been processed
            if (audioElement.hasAttribute('data-processed')) {{
                return;
            }}
            
            // Mark as processed to avoid duplicate event listeners
            audioElement.setAttribute('data-processed', 'true');
            
            // Stop all other audio players when this one starts
            audioElement.addEventListener('play', function() {{
                const allAudios = document.querySelectorAll('audio');
                allAudios.forEach(function(audio) {{
                    if (audio.id !== 'audio_{message_id}' && !audio.paused) {{
                        audio.pause();
                        audio.currentTime = 0;
                    }}
                }});
            }});
            
            // Handle autoplay for the latest message only
            if ({str(should_autoplay).lower()}) {{
                // Stop all currently playing audio first
                const allAudios = document.querySelectorAll('audio');
                allAudios.forEach(function(audio) {{
                    if (!audio.paused) {{
                        audio.pause();
                        audio.currentTime = 0;
                    }}
                }});
                
                // Small delay to ensure other audio is stopped
                setTimeout(function() {{
                    audioElement.play().catch(function(error) {{
                        console.log('Autoplay prevented:', error);
                    }});
                }}, 100);
            }}
        }})();
    </script>
    """
    return audio_html

def create_video_player(video_url, message_id, should_autoplay=False):
    """Create an HTML video player for the avatar response"""
    autoplay_attr = "autoplay muted" if should_autoplay else ""
    
    video_html = f"""
    <div id="video_container_{message_id}" style="margin: 10px 0;">
        <video id="video_{message_id}" controls {autoplay_attr} style="width: 100%; max-width: 400px; border-radius: 10px;">
            <source src="{video_url}" type="video/mp4">
            Your browser does not support the video element.
        </video>
    </div>
    <script>
        (function() {{
            const videoElement = document.getElementById('video_{message_id}');
            
            if (videoElement.hasAttribute('data-processed')) {{
                return;
            }}
            
            videoElement.setAttribute('data-processed', 'true');
            
            // Stop all other videos when this one starts
            videoElement.addEventListener('play', function() {{
                const allVideos = document.querySelectorAll('video');
                const allAudios = document.querySelectorAll('audio');
                
                allVideos.forEach(function(video) {{
                    if (video.id !== 'video_{message_id}' && !video.paused) {{
                        video.pause();
                        video.currentTime = 0;
                    }}
                }});
                
                allAudios.forEach(function(audio) {{
                    if (!audio.paused) {{
                        audio.pause();
                        audio.currentTime = 0;
                    }}
                }});
            }});
            
            if ({str(should_autoplay).lower()}) {{
                setTimeout(function() {{
                    videoElement.play().catch(function(error) {{
                        console.log('Video autoplay prevented:', error);
                    }});
                }}, 100);
            }}
        }})();
    </script>
    """
    return video_html

def create_voice_input_interface():
    """Create voice input interface using Streamlit's built-in audio input"""
    st.markdown("### 🎤 Voice Input")
    
    # Use Streamlit's audio input widget
    audio_bytes = st.audio_input("Record your question")
    
    return audio_bytes

def process_audio_to_text(audio_input):
    """Convert audio to text using speech recognition"""
    if not has_speech_recognition:
        st.error("Speech recognition library not available. Please install speech_recognition.")
        return None
    
    if audio_input is None:
        return None
    
    try:
        # Initialize speech recognizer
        r = sr.Recognizer()
        
        # Handle both UploadedFile and bytes objects
        if hasattr(audio_input, 'getvalue'):
            # It's an UploadedFile object from Streamlit
            audio_bytes = audio_input.getvalue()
        else:
            # It's already bytes
            audio_bytes = audio_input
        
        # Create temporary audio file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.wav') as tmp_audio:
            tmp_audio.write(audio_bytes)
            audio_file_path = tmp_audio.name
        
        try:
            # Load audio file and convert to text
            with sr.AudioFile(audio_file_path) as source:
                # Adjust for ambient noise and record audio
                r.adjust_for_ambient_noise(source, duration=0.5)
                audio = r.record(source)
                
            # Use Google's speech recognition
            text = r.recognize_google(audio)
            logger.info(f"Successfully converted audio to text: {text}")
            return text
            
        except sr.UnknownValueError:
            st.error("Could not understand the audio. Please try speaking more clearly.")
            logger.warning("Speech recognition could not understand audio")
            return None
        except sr.RequestError as e:
            st.error(f"Could not request results from speech recognition service: {e}")
            logger.error(f"Speech recognition service error: {e}")
            return None
        finally:
            # Clean up temporary file
            try:
                os.unlink(audio_file_path)
                logger.info("Temporary audio file deleted")
            except Exception as cleanup_error:
                logger.warning(f"Could not delete temporary audio file: {cleanup_error}")
                
    except Exception as e:
        st.error(f"Error processing audio: {str(e)}")
        logger.error(f"Error in process_audio_to_text: {str(e)}")
        return None

# Avatar settings UI function
def create_avatar_settings_ui():
    """Create enhanced avatar settings UI with custom photo upload"""
    st.header("🎭 Avatar Settings")
    avatar_mode = st.toggle("Enable Avatar Mode", value=st.session_state.get('avatar_mode', False))
    st.session_state.avatar_mode = avatar_mode
    
    if avatar_mode and st.session_state.did_handler.api_key:
        # Avatar type selection
        avatar_type = st.radio(
            "Choose Avatar Type:",
            ["Use Default Presenters", "Upload Your Own Photo"],
            index=0 if not st.session_state.get('use_custom_avatar', False) else 1
        )
        
        st.session_state.use_custom_avatar = (avatar_type == "Upload Your Own Photo")
        
        if st.session_state.use_custom_avatar:
            st.subheader("📸 Upload Your Photo")
            
            # Photo upload
            uploaded_photo = st.file_uploader(
                "Choose your photo",
                type=["jpg", "jpeg", "png"],
                help="Upload a clear, front-facing photo of yourself. Best results with good lighting and neutral background."
            )
            
            if uploaded_photo:
                # Display the uploaded photo
                st.image(uploaded_photo, caption="Your uploaded photo", width=200)
                
                # Upload to D-ID and store URL
                if st.button("🚀 Upload Photo to D-ID"):
                    with st.spinner("Uploading your photo to D-ID..."):
                        image_url = st.session_state.did_handler.upload_custom_image(uploaded_photo)
                        if image_url:
                            st.session_state.custom_avatar_url = image_url
                            st.success("✅ Photo uploaded successfully! Your custom avatar is ready.")
                            st.info(f"🔗 Image URL: {image_url}")
                        else:
                            st.error("❌ Failed to upload photo. Please check your D-ID API key and try again.")
                
                # Show upload status
                if st.session_state.get('custom_avatar_url'):
                    st.success("✅ Custom avatar ready!")
                    if st.button("🗑️ Remove Custom Avatar"):
                        del st.session_state.custom_avatar_url
                        st.rerun()
                else:
                    st.info("👆 Click 'Upload Photo to D-ID' to prepare your custom avatar")
            
            # Photo guidelines
            with st.expander("📋 Photo Guidelines for Best Results"):
                st.markdown("""
                **For optimal avatar quality:**
                
                ✅ **Good Photos:**
                - Clear, front-facing photo
                - Good lighting (natural light preferred)
                - Neutral or plain background
                - Face clearly visible
                - High resolution (at least 512x512px)
                - Person looking directly at camera
                
                ❌ **Avoid:**
                - Blurry or low-quality images
                - Side profiles or angled faces
                - Sunglasses or face coverings
                - Busy or cluttered backgrounds
                - Multiple people in the photo
                - Very dark or overexposed photos
                """)
        
        else:
            # Default presenter selection
            st.subheader("👤 Select Default Presenter")
            presenters = st.session_state.did_handler.get_available_presenters()
            if presenters:
                selected_presenter = st.selectbox(
                    "Select Avatar",
                    options=[presenter[0] for presenter in presenters],
                    format_func=lambda x: next((presenter[1] for presenter in presenters if presenter[0] == x), x),
                    index=0
                )
                st.session_state.did_handler.presenter_id = selected_presenter
        
        st.info("🎭 Avatar will speak responses using D-ID technology")
        
    elif avatar_mode and not st.session_state.did_handler.api_key:
        st.warning("⚠️ D-ID API key not found. Avatar mode disabled.")
        st.info("Add your D-ID API key to secrets to enable avatar features.")

def process_avatar_response(response, user_query):
    """Process avatar response with custom or default avatar"""
    video_url = None
    
    if st.session_state.get('avatar_mode', False) and st.session_state.did_handler.api_key:
        with st.spinner("🎭 Creating avatar response..."):
            logger.info("Avatar mode enabled, creating D-ID talk...")
            
            # Check if using custom avatar
            if st.session_state.get('use_custom_avatar', False) and st.session_state.get('custom_avatar_url'):
                logger.info(f"Using custom avatar with URL: {st.session_state.custom_avatar_url}")
                talk_id = st.session_state.did_handler.create_talk(response, st.session_state.custom_avatar_url)
            else:
                logger.info(f"Using default presenter ID: {st.session_state.did_handler.presenter_id}")
                talk_id = st.session_state.did_handler.create_talk(response)
            
            if talk_id:
                logger.info(f"D-ID talk created with ID: {talk_id}")
                video_url = st.session_state.did_handler.get_talk_url(talk_id)
                logger.info(f"D-ID video URL retrieved: {video_url}")
                if video_url:
                    logger.info("Successfully generated avatar video")
                else:
                    logger.error("Failed to get video URL from D-ID")
            else:
                logger.error("Failed to create D-ID talk")
                st.error("Failed to create avatar response. Please check your D-ID API key and settings.")
    
    return video_url

# Streamlit UI
st.set_page_config(page_title="AI Knowledge-Base Chatbot with Voice", layout="wide")
st.title("🎤 AI Knowledge-Base Chatbot with Voice")
st.write("Upload a PDF and chat with it using voice or text!")

# Initialize session state
if "chatbot" not in st.session_state:
    try:
        st.session_state.chatbot = PDFChatbot()
    except Exception as e:
        logger.error(f"Error initializing chatbot: {str(e)}")
        st.error(f"Error initializing chatbot: {str(e)}")
        st.session_state.chatbot = None

if "voice_handler" not in st.session_state:
    st.session_state.voice_handler = VoiceHandler()

if "conversation" not in st.session_state:
    st.session_state.conversation = None
if "document_processed" not in st.session_state:
    st.session_state.document_processed = False
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "voice_enabled" not in st.session_state:
    st.session_state.voice_enabled = False
if "last_played_audio" not in st.session_state:
    st.session_state.last_played_audio = None
if "message_counter" not in st.session_state:
    st.session_state.message_counter = 0
if "last_autoplay_message_id" not in st.session_state:
    st.session_state.last_autoplay_message_id = -1
if "did_handler" not in st.session_state:
    st.session_state.did_handler = DIDHandler()

# Add auto_play to session state initialization
if "auto_play" not in st.session_state:
    st.session_state.auto_play = True

# Define the layout
left_col, right_col = st.columns([1, 2])

# Sidebar content
with left_col:
    st.header("📄 Upload Document")
    uploaded_file = st.file_uploader("Choose a PDF, TXT, DOCX, PPTX, or Excel file", type=["pdf", "txt", "docx", "pptx", "xlsx"])

    # Voice settings
    st.header("🎤 Voice Settings")
    voice_enabled = st.toggle("Enable Voice Features", value=st.session_state.voice_enabled)
    st.session_state.voice_enabled = voice_enabled
    
    if voice_enabled and st.session_state.voice_handler.api_key:
        # Voice selection
        voices = st.session_state.voice_handler.get_available_voices()
        if voices:
            selected_voice = st.selectbox(
                "Select Voice",
                options=[voice[0] for voice in voices],
                format_func=lambda x: next((voice[1] for voice in voices if voice[0] == x), x),
                index=0
            )
            st.session_state.voice_handler.voice_id = selected_voice
        
        # Voice settings
        st.subheader("🔊 Voice Controls")
        # Add auto_play to session state
        st.session_state.auto_play = st.checkbox("Auto-play responses", 
                                               value=st.session_state.get('auto_play', True))

# Avatar settings
    create_avatar_settings_ui()

    debug_mode = st.toggle("Debug Mode", value=False)

    if "LANGCHAIN_API_KEY" in st.secrets:
        langsmith_url = "https://smith.langchain.com/projects/" + st.secrets.get("LANGCHAIN_PROJECT", "pdf-chatbot")
        st.markdown(f"[View Traces in LangSmith]({langsmith_url})")

    if uploaded_file and not st.session_state.document_processed and st.session_state.chatbot:
        with st.spinner("Processing document..."):
            status_container = st.empty()

            try:
                status_container.info("Extracting text from document...")
                
                # Check file type and process accordingly
                if uploaded_file.name.lower().endswith('.pdf'):
                    documents = st.session_state.chatbot.extract_text_from_pdf(uploaded_file)
                    file_type = "PDF"
                elif uploaded_file.name.lower().endswith('.txt'):
                    documents = st.session_state.chatbot.extract_text_from_txt(uploaded_file)
                    file_type = "TXT"
                elif uploaded_file.name.lower().endswith('.docx'):
                    documents = st.session_state.chatbot.extract_text_from_docx(uploaded_file)
                    file_type = "DOCX"
                elif uploaded_file.name.lower().endswith('.pptx'):
                    documents = st.session_state.chatbot.extract_text_from_pptx(uploaded_file)
                    file_type = "PPTX"
                elif uploaded_file.name.lower().endswith('.xlsx'):
                    documents = st.session_state.chatbot.extract_text_from_excel(uploaded_file)
                    file_type = "Excel"
                else:
                    status_container.error("Unsupported file type. Please upload a PDF, TXT, DOCX, PPTX, or Excel file.")
                    documents = []

                if not documents:
                    status_container.error(f"Failed to extract text from {file_type}. Please try another file.")
                else:
                    if file_type == "PDF":
                        status_container.success(f"✅ Extracted {len(documents)} pages from PDF")
                    elif file_type == "TXT":
                        status_container.success(f"✅ Successfully loaded TXT file")
                    elif file_type == "DOCX":
                        status_container.success(f"✅ Successfully loaded DOCX file")
                    elif file_type == "PPTX":
                        status_container.success(f"✅ Extracted {len(documents)} slides from PPTX")
                    elif file_type == "Excel":
                        status_container.success(f"✅ Extracted {len(documents)} sheets from Excel file")

                    status_container.info("Chunking text...")
                    chunks = st.session_state.chatbot.chunk_text(documents)
                    status_container.success(f"✅ Created {len(chunks)} text chunks")

                    namespace = uploaded_file.name.replace(" ", "_").lower()

                    status_container.info("Creating embeddings and storing in Pinecone...")
                    vectorstore = st.session_state.chatbot.create_embeddings_and_store(chunks, namespace)
                    status_container.success("✅ Created embeddings and stored in Pinecone")

                    status_container.info("Setting up conversation chain...")
                    st.session_state.conversation = st.session_state.chatbot.get_conversational_chain(vectorstore)
                    st.session_state.document_processed = True
                    status_container.empty()
                    st.success("Document processed successfully!")
            except Exception as e:
                logger.error(f"Error processing document: {str(e)}")
                status_container.error(f"Error processing document: {str(e)}")

    if st.session_state.document_processed:
        st.subheader("📋 Document Information")
        st.info(f"📁 Filename: {uploaded_file.name}")

        if st.button("🔄 Process Another Document"):
            st.session_state.document_processed = False
            st.session_state.conversation = None
            st.session_state.chat_history = []
            st.session_state.message_counter = 0
            st.session_state.last_played_audio = None
            st.session_state.last_autoplay_message_id = -1
            st.rerun()

# Chat interface
with right_col:
    st.header("💬 Chat with your Knowledge-Base")

    if not st.session_state.document_processed:
        st.info("Please upload a PDF document to start chatting.")
    else:
        # Display chat history with proper audio/video control
        for i, message in enumerate(st.session_state.chat_history):
            if message["role"] == "user":
                st.chat_message("user").write(message["content"])
            else:
                st.chat_message("assistant").write(message["content"])
                
                message_id = message.get("message_id", i)
                is_latest = (i == len(st.session_state.chat_history) - 1)
                
                # Show avatar video if available
                if st.session_state.get('avatar_mode', False) and "video_url" in message:
                    should_autoplay = (
                        is_latest and 
                        st.session_state.auto_play and 
                        st.session_state.get("last_autoplay_message_id", -1) != message_id
                    )
                    
                    if should_autoplay:
                        st.session_state.last_autoplay_message_id = message_id
                    
                    st.markdown(
                        create_video_player(message["video_url"], message_id, should_autoplay), 
                        unsafe_allow_html=True
                    )
                    
                # Show audio player if voice is enabled and no avatar mode
                elif st.session_state.voice_enabled and "audio" in message and not st.session_state.get('avatar_mode', False):
                    should_autoplay = (
                        is_latest and 
                        st.session_state.auto_play and 
                        st.session_state.get("last_autoplay_message_id", -1) != message_id
                    )
                    
                    if should_autoplay:
                        st.session_state.last_autoplay_message_id = message_id
                    
                    st.markdown(
                        create_audio_player(message["audio"], message_id, should_autoplay), 
                        unsafe_allow_html=True
                    )

        # Voice input section
        user_query = None
        
        if st.session_state.voice_enabled and has_speech_recognition:
            # Voice input using Streamlit's audio input
            audio_input = create_voice_input_interface()
            
            if audio_input is not None:
                with st.spinner("🎤 Converting speech to text..."):
                    voice_text = process_audio_to_text(audio_input)
                    if voice_text:
                        st.success(f"🎤 Voice input: {voice_text}")
                        user_query = voice_text

        # Text input (always available)
        if not user_query:
            user_query = st.chat_input("Ask a question about your document")

        # Process the query
        if user_query:
            # Add a global audio stop before processing new query
            st.markdown("""
            <script>
                (function stopAllAudio() {
                    const allAudios = document.querySelectorAll('audio');
                    allAudios.forEach(function(audio) {
                        audio.pause();
                        audio.currentTime = 0;
                    });
                })();
            </script>
            """, unsafe_allow_html=True)
            
            st.chat_message("user").write(user_query)
            st.session_state.chat_history.append({"role": "user", "content": user_query})

            with st.spinner("🤔 Thinking..."):
                try:
                    response, sources = st.session_state.chatbot.process_query(
                        st.session_state.conversation, user_query
                    )
                    
                    # Generate avatar video or audio based on settings
                    audio_content = None
                    video_url = process_avatar_response(response, user_query)
                    
                except Exception as e:
                    logger.error(f"Error processing query: {str(e)}")
                    response = f"I'm sorry, I encountered an error while processing your query: {str(e)}"
                    sources = []
                    audio_content = None
                    video_url = None

            # Display response
            st.chat_message("assistant").write(response)
            
            # Store response in chat history with unique message ID
            st.session_state.message_counter += 1
            message_data = {
                "role": "assistant", 
                "content": response,
                "message_id": st.session_state.message_counter
            }
            
            # Add audio or video data to message
            if video_url:
                message_data["video_url"] = video_url
            elif audio_content:
                message_data["audio"] = audio_content
                
            st.session_state.chat_history.append(message_data)

            # Display avatar video or play audio
            if video_url:
                st.markdown("🎭 **Avatar Response:**")
                st.markdown(
                    create_video_player(
                        video_url, 
                        st.session_state.message_counter, 
                        st.session_state.auto_play  # Use session state variable
                    ), 
                    unsafe_allow_html=True
                )
                
                if st.session_state.auto_play:  # Use session state variable
                    st.session_state.last_autoplay_message_id = st.session_state.message_counter
                    
            elif audio_content:
                st.markdown("🔊 **Audio Response:**")
                st.markdown(
                    create_audio_player(
                        audio_content, 
                        st.session_state.message_counter, 
                        st.session_state.auto_play  # Use session state variable
                    ), 
                    unsafe_allow_html=True
                )
                
                if st.session_state.auto_play:  # Use session state variable
                    st.session_state.last_autoplay_message_id = st.session_state.message_counter

            # Show sources
            if sources:
                with st.expander("📚 Sources"):
                    for i, source in enumerate(sources):
                        st.write(f"**Source {i + 1}:**")
                        st.write(source.page_content)
                        st.write(f"📄 Page: {source.metadata.get('page', 'N/A')}")
                        st.divider()

            # Add this after processing the response:
            if debug_mode:
                st.info("Debug Information:")
                st.json({
                    "Avatar Mode": st.session_state.get('avatar_mode', False),
                    "D-ID API Key Present": bool(st.session_state.did_handler.api_key),
                    "Presenter ID": st.session_state.did_handler.presenter_id,
                    "Response Length": len(response),
                    "Video URL": video_url if video_url else "None",
                    "Auto Play": st.session_state.auto_play,
                    "Message Counter": st.session_state.message_counter
                })

# Setup instructions
if not st.session_state.document_processed:
    with st.expander("🔧 How to set up API keys and dependencies"):
        st.markdown("""
        ### Setting up your Streamlit Cloud secrets

        In your Streamlit Cloud dashboard:
        1. Go to your app settings
        2. Navigate to the "Secrets" section
        3. Add the following secrets:

        ```toml
        # Required API keys
        PINECONE_API_KEY = "your-pinecone-api-key"
        GOOGLE_API_KEY = "your-google-api-key"

        # Required for Pinecone client v2.x.x
        PINECONE_ENVIRONMENT = "gcp-starter"  # Or your specific environment

        # Voice functionality (ElevenLabs)
        ELEVENLABS_API_KEY = "your-elevenlabs-api-key"
        ELEVENLABS_VOICE_ID = "21m00Tcm4TlvDq8ikWAM"  # Optional: specific voice ID

        # Avatar functionality (D-ID)
        DID_API_KEY = "your-did-api-key"
        DID_PRESENTER_ID = "your-presenter-id"  # Optional: specific presenter ID

        # Optional LangSmith configuration (for advanced monitoring)
        LANGCHAIN_API_KEY = "your-langsmith-api-key"
        LANGCHAIN_PROJECT = "pdf-chatbot"
        LANGCHAIN_TRACING_V2 = "true"
        ```

        ### Requirements.txt for Streamlit Cloud
        
        Add this to your requirements.txt file:
        ```
        streamlit
        langchain
        langchain-community
        langchain-huggingface
        langchain-google-genai
        langchain-pinecone
        pinecone-client
        google-generativeai
        PyMuPDF
        sentence-transformers
        requests
        SpeechRecognition
        ```

        ### 🎤 Voice Features Setup
        
        1. **ElevenLabs Account**: Sign up at [ElevenLabs.io](https://elevenlabs.io)
        2. **Get API Key**: Go to your profile settings and copy your API key
        3. **Choose Voice**: Browse available voices and copy the voice ID (optional)
        4. **Browser Support**: Voice input uses Streamlit's built-in audio recording
        
        ### 🔊 Voice Features Include:
        - **Voice Input**: Record your questions using Streamlit's audio input
        - **Speech-to-Text**: Automatic transcription using Google Speech Recognition
        - **AI Processing**: Questions sent to Google Gemini for processing
        - **Text-to-Speech**: Responses converted to speech using ElevenLabs
        - **Audio Playback**: Seamless audio responses in the chat interface

        ### 🎭 Avatar Features Setup

        1. **D-ID Account**: Sign up at [D-ID.com](https://www.d-id.com/)
        2. **Get API Key**: Go to your account settings and copy your API key
        3. **Choose Presenter**: Browse available presenters and copy the presenter ID (optional)
        4. **Browser Support**: Avatar videos are generated using D-ID's API and played in the browser
        """)

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; opacity: 0.7; font-size: 0.9rem;">
    Built with ❤️ using <strong>Streamlit</strong>, <strong>LangChain</strong>, <strong>Pinecone</strong>, <strong>Google Gemini</strong>, and <strong>ElevenLabs</strong>
    <br>
    🎤 Voice-enabled AI PDF Chatbot
</div>
""", unsafe_allow_html=True)
